import os
import json
import shutil
import zipfile
from pathlib import Path
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

BASE_DIR = Path("F:/0workbuddy文件存放/2026-08-31-00-02-43")
MATERIAL_DIR = Path("C:/Users/DELL/Desktop/我的资料")
OUTPUT_DIR = BASE_DIR / "portfolio_assets"
OUTPUT_DIR.mkdir(exist_ok=True)

# 1. Extract resume PDF text
resume_pdf = MATERIAL_DIR / "长安大学-张瑞娟-硕士-27届.pdf"
resume_text_path = OUTPUT_DIR / "resume_text.txt"

print(f"Extracting: {resume_pdf.name}")
with pdfplumber.open(resume_pdf) as pdf, open(resume_text_path, "w", encoding="utf-8") as f:
    for i, page in enumerate(pdf.pages):
        text = page.extract_text() or ""
        f.write(f"\n--- Page {i+1} ---\n")
        f.write(text)
        f.write("\n")

# 2. Extract portfolio PDF text (and render pages as images if not too many)
portfolio_pdf = MATERIAL_DIR / "张瑞娟＋作品集.pdf"
portfolio_text_path = OUTPUT_DIR / "portfolio_text.txt"
portfolio_images_dir = OUTPUT_DIR / "portfolio_pages"
portfolio_images_dir.mkdir(exist_ok=True)

print(f"Extracting: {portfolio_pdf.name}")
with pdfplumber.open(portfolio_pdf) as pdf, open(portfolio_text_path, "w", encoding="utf-8") as f:
    for i, page in enumerate(pdf.pages):
        text = page.extract_text() or ""
        f.write(f"\n--- Page {i+1} ---\n")
        f.write(text)
        f.write("\n")
        # Render first 15 pages as images
        if i < 15:
            img = page.to_image(resolution=150)
            img.save(portfolio_images_dir / f"page_{i+1:02d}.png")

# 3. Extract PPT text and images
ppt_path = MATERIAL_DIR / "优秀实习生答辩·张瑞娟.pptx"
ppt_images_dir = OUTPUT_DIR / "ppt_images"
ppt_images_dir.mkdir(exist_ok=True)
ppt_text_path = OUTPUT_DIR / "ppt_text.txt"

print(f"Extracting: {ppt_path.name}")
prs = Presentation(ppt_path)
with open(ppt_text_path, "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                f.write(shape.text.strip() + "\n")

# Extract images from PPTX by unzipping
pptx_zip_dir = OUTPUT_DIR / "pptx_unzipped"
pptx_zip_dir.mkdir(exist_ok=True)
with zipfile.ZipFile(ppt_path, "r") as zip_ref:
    zip_ref.extractall(pptx_zip_dir)

media_dir = pptx_zip_dir / "ppt" / "media"
if media_dir.exists():
    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
    video_exts = {".mp4", ".mov", ".avi", ".wmv", ".mkv", ".webm"}
    img_idx = 1
    vid_idx = 1
    for item in sorted(media_dir.iterdir()):
        if item.suffix.lower() in image_exts:
            target = ppt_images_dir / f"ppt_image_{img_idx:03d}{item.suffix.lower()}"
            shutil.copy2(item, target)
            img_idx += 1
        elif item.suffix.lower() in video_exts:
            target = OUTPUT_DIR / f"ppt_video_{vid_idx:03d}{item.suffix.lower()}"
            shutil.copy2(item, target)
            vid_idx += 1

print("Extraction complete.")
print(f"Output directory: {OUTPUT_DIR}")
